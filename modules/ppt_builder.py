import os
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import config

def hex_to_rgb(hex_str):
    """
    Converts a hex color string (e.g., '#1B365D') to an RGBColor object.
    """
    hex_str = hex_str.lstrip('#')
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)

def get_theme_palette(theme):
    """
    Returns hex colors and font styling guidelines for pptx builders.
    """
    t = theme.lower()
    if t == "corporate":
        return {
            "bg": "#F4F6F9",
            "header": "#1B365D",      # Navy
            "accent": "#4A777A",      # Slate Teal
            "highlight": "#D99B26",   # Gold
            "text": "#333333",
            "card_bg": "#FFFFFF",
            "font_title": "Arial",
            "font_body": "Calibri"
        }
    elif t == "consulting":
        return {
            "bg": "#FAF9F6",          # Warm Cream
            "header": "#2D5A27",      # Sage Green
            "accent": "#4F5D75",      # Slate Gray
            "highlight": "#D1A153",   # Warm Ochre
            "text": "#2C2C2C",
            "card_bg": "#FFFFFF",
            "font_title": "Georgia",
            "font_body": "Calibri"
        }
    else: # minimal (default)
        return {
            "bg": "#FFFFFF",          # Widescreen White
            "header": "#222222",      # Charcoal
            "accent": "#777777",      # Cool Gray
            "highlight": "#007ACC",   # Tech Blue
            "text": "#111111",
            "card_bg": "#F8F9FA",
            "font_title": "Helvetica",
            "font_body": "Arial"
        }

def build_presentation(slides_outline, df, chart_paths, theme="minimal", target_audience="CEO"):
    """
    Builds a professional PowerPoint deck based on the outline structure and saved charts.
    Returns:
        str: Absolute filepath of the generated presentation (.pptx)
    """
    palette = get_theme_palette(theme)
    
    # 1. Initialize Presentation with Widescreen (16:9) dimensions
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # completely blank layout
    
    slides = slides_outline.get("slides", [])
    
    for idx, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        
        # Add slide background color
        _set_slide_background(slide, palette["bg"])
        
        # Check if first slide (Title slide)
        if idx == 0:
            _build_title_slide(slide, slide_data, palette)
            continue
            
        # Check if second slide (Executive Summary)
        if idx == 1:
            _build_exec_summary_slide(slide, slide_data, palette)
            continue
            
        # Check if slide has a chart
        visual_type = slide_data.get("visual_required", "None")
        chart_path = chart_paths.get(idx)
        
        if chart_path and os.path.exists(chart_path) and visual_type != "None":
            _build_split_chart_slide(slide, slide_data, chart_path, palette)
        else:
            # Check if this is the final slide (typically recommendations)
            if idx == len(slides) - 1:
                _build_recommendations_slide(slide, slide_data, palette)
            else:
                _build_standard_text_slide(slide, slide_data, palette)
                
    # Save presentation
    filename = f"ExecuDeck_{theme}_{target_audience}_{uuid.uuid4().hex[:6]}.pptx"
    filepath = config.OUTPUT_PRES_DIR / filename
    prs.save(filepath)
    
    return str(filepath)

def _set_slide_background(slide, bg_hex):
    """
    Draws a full-slide rectangle as the background color.
    """
    left = top = 0
    width = Inches(13.333)
    height = Inches(7.5)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = hex_to_rgb(bg_hex)
    rect.line.fill.background() # no outline

def _add_slide_header(slide, title, purpose, palette):
    """
    Adds a standardized top title and subtitle/purpose to the slide.
    """
    # Header Textbox
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    # Title
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = palette["font_title"]
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = hex_to_rgb(palette["header"])
    p_title.space_after = Pt(4)
    
    # Subtitle / Purpose
    p_sub = tf.add_paragraph()
    p_sub.text = purpose
    p_sub.font.name = palette["font_body"]
    p_sub.font.size = Pt(12)
    p_sub.font.italic = True
    p_sub.font.color.rgb = hex_to_rgb(palette["accent"])
    
    # Bottom separator line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.6), Inches(11.83), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(palette["highlight"])
    line.line.fill.background()

def _build_title_slide(slide, data, palette):
    """
    Builds a striking, McKinsey-style Title Slide.
    """
    # Draw left decorative block (1/3 of slide)
    left_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
    left_block.fill.solid()
    left_block.fill.fore_color.rgb = hex_to_rgb(palette["header"])
    left_block.line.fill.background()
    
    # Add highlight strip on the edge
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.35), Inches(0), Inches(0.15), Inches(7.5))
    strip.fill.solid()
    strip.fill.fore_color.rgb = hex_to_rgb(palette["highlight"])
    strip.line.fill.background()
    
    # Title text box
    title_box = slide.shapes.add_textbox(Inches(5.0), Inches(2.2), Inches(7.5), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    # Title
    p_title = tf.paragraphs[0]
    p_title.text = data.get("title", "ExecuDeck Presentation")
    p_title.font.name = palette["font_title"]
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = hex_to_rgb(palette["header"])
    p_title.space_after = Pt(12)
    
    # Subtitle
    p_sub = tf.add_paragraph()
    p_sub.text = data.get("purpose", "Executive Presentation Generator")
    p_sub.font.name = palette["font_body"]
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = hex_to_rgb(palette["accent"])
    p_sub.space_after = Pt(24)
    
    # Metadata info (Enterprise local statement)
    p_meta = tf.add_paragraph()
    p_meta.text = "Prepared by ExecuDeck AI  |  Secure Local Inference Environment"
    p_meta.font.name = palette["font_body"]
    p_meta.font.size = Pt(11)
    p_meta.font.color.rgb = hex_to_rgb(palette["accent"])
    p_meta.font.bold = True

def _build_exec_summary_slide(slide, data, palette):
    """
    Builds a two-column Executive Summary & KPI cards slide.
    """
    _add_slide_header(slide, data.get("title", "Executive Context"), data.get("purpose", ""), palette)
    
    # Left column: Narrative Card
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(2.0), Inches(5.6), Inches(4.7))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = hex_to_rgb(palette["card_bg"])
    left_card.line.color.rgb = hex_to_rgb(palette["accent"])
    left_card.line.width = Pt(1)
    
    # Left card text box
    left_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.1), Inches(4.3))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_right = tf_left.margin_bottom = 0
    
    p_lh = tf_left.paragraphs[0]
    p_lh.text = "Executive Strategic Analysis"
    p_lh.font.name = palette["font_title"]
    p_lh.font.size = Pt(18)
    p_lh.font.bold = True
    p_lh.font.color.rgb = hex_to_rgb(palette["header"])
    p_lh.space_after = Pt(14)
    
    # Add bullet points
    for bullet in data.get("bullet_points", []):
        p = tf_left.add_paragraph()
        p.text = "• " + bullet
        p.font.name = palette["font_body"]
        p.font.size = Pt(13)
        p.font.color.rgb = hex_to_rgb(palette["text"])
        p.space_after = Pt(10)
        
    # Right column: Metric Cards
    # We display up to 3 metrics cards vertically
    right_box_left = Inches(6.8)
    card_width = Inches(5.78)
    card_height = Inches(1.35)
    gap = Inches(0.3)
    
    # If no metrics, make up standard summary points
    metrics = data.get("bullet_points", []) # Default if metrics not distinct
    
    for i in range(min(3, len(metrics))):
        top_pos = Inches(2.0) + i * (card_height + gap)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_box_left, top_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(palette["card_bg"])
        card.line.color.rgb = hex_to_rgb(palette["highlight"])
        card.line.width = Pt(1.5)
        
        # Card Text Box
        card_box = slide.shapes.add_textbox(right_box_left + Inches(0.2), top_pos + Inches(0.15), card_width - Inches(0.4), card_height - Inches(0.3))
        tf_card = card_box.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = tf_card.margin_top = tf_card.margin_right = tf_card.margin_bottom = 0
        
        # Split text into KPI title and description if contains ':'
        text = metrics[i]
        title_text = "Metric Update"
        desc_text = text
        if ":" in text:
            parts = text.split(":", 1)
            title_text = parts[0].strip()
            desc_text = parts[1].strip()
            
        p_ct = tf_card.paragraphs[0]
        p_ct.text = title_text
        p_ct.font.name = palette["font_body"]
        p_ct.font.size = Pt(11)
        p_ct.font.bold = True
        p_ct.font.color.rgb = hex_to_rgb(palette["accent"])
        p_ct.space_after = Pt(3)
        
        p_cd = tf_card.add_paragraph()
        p_cd.text = desc_text
        p_cd.font.name = palette["font_title"]
        p_cd.font.size = Pt(15)
        p_cd.font.bold = True
        p_cd.font.color.rgb = hex_to_rgb(palette["header"])

def _build_standard_text_slide(slide, data, palette):
    """
    Builds a standard information slide with bullet points.
    """
    _add_slide_header(slide, data.get("title", "Insight"), data.get("purpose", ""), palette)
    
    # Bullet points text frame
    bullets_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.83), Inches(4.5))
    tf = bullets_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    # Loop over bullets
    bullets = data.get("bullet_points", [])
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = "•  " + bullet
        p.font.name = palette["font_body"]
        # Size guard: adjust font size if there are too many bullets
        p.font.size = Pt(16) if len(bullets) <= 4 else Pt(14)
        p.font.color.rgb = hex_to_rgb(palette["text"])
        p.space_after = Pt(14)

def _build_split_chart_slide(slide, data, chart_path, palette):
    """
    Builds a slide containing a narrative column on the left and the chart image on the right.
    """
    _add_slide_header(slide, data.get("title", "Data Visualization"), data.get("purpose", ""), palette)
    
    # Left narrative block
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.1), Inches(5.8), Inches(4.6))
    tf = left_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    bullets = data.get("bullet_points", [])
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = "•  " + bullet
        p.font.name = palette["font_body"]
        # Size guard
        p.font.size = Pt(13.5) if len(bullets) <= 4 else Pt(12)
        p.font.color.rgb = hex_to_rgb(palette["text"])
        p.space_after = Pt(10)
        
    # Right chart insertion
    # Auto-centered inside the right 6.5 inch block
    chart_left = Inches(7.0)
    chart_top = Inches(2.0)
    chart_width = Inches(5.6)
    chart_height = Inches(4.5)
    
    slide.shapes.add_picture(chart_path, chart_left, chart_top, width=chart_width, height=chart_height)

def _build_recommendations_slide(slide, data, palette):
    """
    Builds a modern layout of numbered strategic cards.
    """
    _add_slide_header(slide, data.get("title", "Strategic Recommendations"), data.get("purpose", ""), palette)
    
    recommendations = data.get("bullet_points", [])
    num_recs = len(recommendations)
    
    if num_recs == 0:
        return
        
    # Standardize to a 3-column horizontal grid
    col_width = Inches(3.7)
    card_height = Inches(4.4)
    gap = Inches(0.4)
    start_left = Inches(0.75)
    
    # Recalculate margins if more/less recommendations
    if num_recs == 2:
        col_width = Inches(5.6)
        gap = Inches(0.6)
    elif num_recs == 4:
        col_width = Inches(2.7)
        gap = Inches(0.3)
        
    for i in range(min(4, num_recs)):
        current_left = start_left + i * (col_width + gap)
        current_top = Inches(2.1)
        
        # Draw Card Shape
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, current_left, current_top, col_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(palette["card_bg"])
        card.line.color.rgb = hex_to_rgb(palette["highlight"])
        card.line.width = Pt(1.5)
        
        # Number indicator shape (Top edge)
        num_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, current_left + Inches(0.3), current_top + Inches(0.3), Inches(0.75), Inches(0.75))
        num_bg.fill.solid()
        num_bg.fill.fore_color.rgb = hex_to_rgb(palette["header"])
        num_bg.line.fill.background()
        
        tf_num = num_bg.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = str(i + 1)
        p_num.alignment = PP_ALIGN.CENTER
        p_num.font.name = palette["font_title"]
        p_num.font.size = Pt(18)
        p_num.font.bold = True
        p_num.font.color.rgb = hex_to_rgb("#FFFFFF")
        
        # Card Text Frame
        tb = slide.shapes.add_textbox(current_left + Inches(0.3), current_top + Inches(1.3), col_width - Inches(0.6), card_height - Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        text = recommendations[i]
        title_text = f"Action Item {i+1}"
        body_text = text
        
        if "." in text and len(text.split(".", 1)[0]) < 25:
            parts = text.split(".", 1)
            title_text = parts[0].strip()
            body_text = parts[1].strip()
        elif ":" in text and len(text.split(":", 1)[0]) < 25:
            parts = text.split(":", 1)
            title_text = parts[0].strip()
            body_text = parts[1].strip()
            
        p_title = tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = palette["font_title"]
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = hex_to_rgb(palette["header"])
        p_title.space_after = Pt(8)
        
        p_body = tf.add_paragraph()
        p_body.text = body_text
        p_body.font.name = palette["font_body"]
        p_body.font.size = Pt(11)
        p_body.font.color.rgb = hex_to_rgb(palette["text"])
        p_body.space_after = Pt(0)
